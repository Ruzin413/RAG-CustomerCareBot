import os
import json
import logging
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

from transformers import AutoTokenizer
import difflib

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        """Initialize document processor"""
        
        # Load tokenizer for token estimation
        try:
            self.tokenizer = AutoTokenizer.from_pretrained("gpt2")
            logger.info("✓ GPT-2 tokenizer loaded for token estimation")
        except Exception as e:
            self.tokenizer = None
            logger.warning(f"⚠️ Could not load tokenizer: {e}. Using character-based estimation")

    def extract_text_from_docx(self, file_stream):
        """Extract all text from DOCX file"""
        try:
            logger.info(f"📂 Reading DOCX file...")
            doc = Document(file_stream)
            full_text = []
            
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    full_text.append(para.text.strip())
                if (i + 1) % 50 == 0:
                    logger.info(f"  - Read {i+1} paragraphs...")
            
            logger.info(f"✓ Extracted {len(full_text)} non-empty paragraphs from DOCX")
            return full_text
            
        except Exception as e:
            logger.error(f"❌ Error extracting DOCX: {e}")
            raise

    def extract_text_from_pdf(self, file_stream):
        """Extract all text from PDF using PyMuPDF, preserving page structure"""
        try:
            logger.info(f"📂 Reading PDF file content...")
            pdf_bytes = file_stream.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            pages_text = []
            
            total_pages = len(doc)
            logger.info(f"📄 PDF has {total_pages} pages. Starting extraction...")
            
            for page in doc:
                text = page.get_text().strip()
                if text:
                    pages_text.append(text)
                
                if (page.number + 1) % 5 == 0 or (page.number + 1) == total_pages:
                    logger.info(f"  - Extracted page {page.number + 1}/{total_pages}...")
            
            doc.close()
            logger.info(f"✓ Extracted {len(pages_text)} non-empty pages from PDF")
            return pages_text
            
        except Exception as e:
            logger.error(f"❌ Error extracting PDF: {e}")
            raise

    def extract_text_from_ppt(self, file_stream):
        """Extract all text from PowerPoint (.pptx) file"""
        try:
            logger.info(f"📂 Reading PPTX file...")
            prs = Presentation(file_stream)
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    text = "\n".join(slide_content)
                    slides_text.append(text)
                
                if (i + 1) % 10 == 0:
                    logger.info(f"  - Read {i+1} slides...")
            
            logger.info(f"✓ Extracted {len(slides_text)} non-empty slides from PPTX")
            return slides_text
            
        except Exception as e:
            logger.error(f"❌ Error extracting PPTX: {e}")
            raise

    def extract_chunks(self, file_name, text_units, target_tokens=300, overlap_tokens=50, kb_name="General"):
        """
        Groups text units into RAG-friendly chunks using LangChain's RecursiveCharacterTextSplitter.
        This respects natural boundaries like paragraphs and sentences.
        """
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        import time
        import uuid

        doc_id = file_name.replace(" ", "_").lower()
        
        # Combine all units into one full text first to allow LangChain to find optimal splits
        full_text = "\n\n".join(text_units)

        logger.info(f"✂️  Starting LangChain chunking for {file_name} (Chunk Size: 500 chars)...")

        # Initialize RecursiveCharacterTextSplitter
        # tries paragraph → line → sentence → word
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=50,
            separators=["\n\n", "\n", ". ", " "]
        )

        # Split the text
        split_texts = text_splitter.split_text(full_text)

        chunks = []
        for i, text in enumerate(split_texts):
            chunks.append({
                "doc_id": doc_id,
                "chunk_id": f"{doc_id}#{i:03d}",
                "kb_name": kb_name,
                "source": {"file": file_name},
                "text": text,
                "tags": [],
                "verified": True,
                "created_at": time.strftime("%Y-%m-%dT%H:%M:%S")
            })

        logger.info(f"✅ LangChain chunking complete. Generated {len(chunks)} context-aware chunks.")
        return chunks
        return chunks